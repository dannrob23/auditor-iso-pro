#!/usr/bin/env python3
"""Convierte archivos Excel, PDF y DOCX a Markdown para Obsidian."""

import os
import sys
from pathlib import Path

SRC_DIR = Path("/home/dannrob/auditor-iso-pro/knowledge_base")
OUT_DIR = SRC_DIR / "markdown"
OUT_DIR.mkdir(parents=True, exist_ok=True)


def pdf_to_md(pdf_path: Path) -> str:
    import pdfplumber

    with pdfplumber.open(str(pdf_path)) as pdf:
        pages_md = []
        for i, page in enumerate(pdf.pages, 1):
            text = page.extract_text()
            tables = page.extract_tables()
            md_lines = [f"## Página {i}\n"]
            if text:
                md_lines.append(text)
            for table in tables:
                if not table:
                    continue
                md_lines.append("\n")
                md_table = _table_to_md(table)
                md_lines.append(md_table)
            pages_md.append("\n".join(md_lines))

    return "\n\n---\n\n".join(pages_md)


def excel_to_md(excel_path: Path) -> str:
    if excel_path.suffix.lower() == ".xls":
        import xlrd
        wb = xlrd.open_workbook(str(excel_path))
        sheets_md = []
        for sheet in wb.sheets():
            md_lines = [f"## Hoja: {sheet.name}\n"]
            rows = []
            for row_idx in range(sheet.nrows):
                row = [str(sheet.cell_value(row_idx, col_idx) or "") for col_idx in range(sheet.ncols)]
                rows.append(row)
            if rows:
                md_lines.append(_rows_to_md(rows))
            sheets_md.append("\n".join(md_lines))
        return "\n\n---\n\n".join(sheets_md)
    else:
        from openpyxl import load_workbook
        wb = load_workbook(str(excel_path), data_only=True)
        sheets_md = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            md_lines = [f"## Hoja: {sheet_name}\n"]
            rows = []
            for row in ws.iter_rows(values_only=True):
                rows.append([str(cell or "") for cell in row])
            if rows:
                md_lines.append(_rows_to_md(rows))
            sheets_md.append("\n".join(md_lines))
            ws_max = ws.max_row or 0
            if ws_max > 100:
                md_lines.append(f"\n*(Hoja truncada: mostrando {len(rows)} de {ws_max}+ filas)*")
        return "\n\n---\n\n".join(sheets_md)


def docx_to_md(docx_path: Path) -> str:
    from docx import Document

    doc = Document(str(docx_path))
    md_lines = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            md_lines.append("")
            continue
        style = para.style.name if para.style else ""
        if "Heading 1" in style:
            md_lines.append(f"# {text}")
        elif "Heading 2" in style:
            md_lines.append(f"## {text}")
        elif "Heading 3" in style:
            md_lines.append(f"### {text}")
        elif "Heading" in style:
            md_lines.append(f"#### {text}")
        elif "List" in style or text.startswith("-") or text.startswith("•"):
            md_lines.append(f"- {text.lstrip('•- ')}")
        else:
            md_lines.append(text)

    for table in doc.tables:
        md_lines.append("")
        rows = []
        for row in table.rows:
            rows.append([cell.text.strip() for cell in row.cells])
        if rows:
            md_lines.append(_rows_to_md(rows))

    return "\n".join(md_lines)


def pptx_to_md(pptx_path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    md_lines = []

    for i, slide in enumerate(prs.slides, 1):
        md_lines.append(f"## Diapositiva {i}\n")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    md_lines.append(text)
        md_lines.append("\n---\n")

    return "\n".join(md_lines)


def txt_to_md(txt_path: Path) -> str:
    return txt_path.read_text(encoding="utf-8", errors="replace")


def _rows_to_md(rows: list) -> str:
    if not rows:
        return ""
    lines = []
    header = rows[0]
    lines.append("| " + " | ".join(header) + " |")
    lines.append("| " + " | ".join(["---"] * len(header)) + " |")
    for row in rows[1:]:
        padded = row + [""] * (len(header) - len(row))
        lines.append("| " + " | ".join(padded[:len(header)]) + " |")
    return "\n".join(lines)


def _table_to_md(table: list) -> str:
    rows = [[str(c or "") for c in row] for row in table]
    return _rows_to_md(rows)


CONVERTERS = {
    ".pdf": pdf_to_md,
    ".xlsx": excel_to_md,
    ".xls": excel_to_md,
    ".docx": docx_to_md,
    ".pptx": pptx_to_md,
    ".txt": txt_to_md,
}


def main():
    files = sorted(SRC_DIR.iterdir())
    converted = 0
    failed = 0

    for f in files:
        if f.is_dir() or f.suffix.lower() not in CONVERTERS:
            continue

        out_path = OUT_DIR / f"{f.stem}.md"
        print(f"🔄 {f.name} → {out_path.name} ...", end=" ")

        try:
            converter = CONVERTERS[f.suffix.lower()]
            md_content = converter(f)
            out_path.write_text(md_content, encoding="utf-8")
            size_kb = out_path.stat().st_size / 1024
            print(f"✓ ({size_kb:.1f} KB)")
            converted += 1
        except Exception as e:
            print(f"✗ ERROR: {e}")
            failed += 1

    print(f"\n✅ {converted} archivos convertidos, ❌ {failed} fallos")
    print(f"📁 Salida: {OUT_DIR}")


if __name__ == "__main__":
    main()
